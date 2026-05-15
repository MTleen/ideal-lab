#!/usr/bin/env python3
"""
image-export.py - Image mode PPTX export

Usage:
    python3 image-export.py <slide-deck-dir>

Input:
    <slide-deck-dir>/images_output/NN-slide-*.png
    <slide-deck-dir>/notes/total.md (optional, for speaker notes)

Output:
    <slide-deck-dir>/<topic-slug>.pptx
"""

import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("Error: Pillow not installed. Run: pip install Pillow")
    sys.exit(1)


def parse_notes(notes_path: Path) -> dict[str, str]:
    """Parse notes/total.md into a {slide_num: note_text} map."""
    if not notes_path.exists():
        return {}

    content = notes_path.read_text(encoding="utf-8")

    # Split by "# NN" headings. MUST use re.MULTILINE so ^ matches each line.
    sections = re.split(r"^# (\d+)", content, flags=re.MULTILINE)

    notes_map: dict[str, str] = {}
    for i in range(1, len(sections), 2):
        slide_num = sections[i].strip().split(" - ")[0].split(" ")[0]
        note_text = sections[i + 1].strip() if i + 1 < len(sections) else ""

        # Clean up formatting markers
        note_text = re.sub(r"\*\*时长.*?\*\*", "", note_text)
        lines = []
        for line in note_text.split("\n"):
            line = re.sub(r"^[①②③④⑤⑥⑦⑧⑨⑩]\s*", "- ", line)
            line = re.sub(r"^\*\*\[Transition:?\s*(.*?)\]\*\*", r"Transition: \1", line)
            line = re.sub(r"^\*\*\[Pause\]\*\*", "Pause", line)
            lines.append(line)

        cleaned = "\n".join(lines).strip()
        if cleaned:
            notes_map[slide_num.zfill(2)] = cleaned

    return notes_map


def find_topic_slug(deck_dir: Path) -> str:
    """Derive topic slug from directory name."""
    return deck_dir.name


def export_pptx(deck_dir: Path):
    images_dir = deck_dir / "images_output"
    notes_path = deck_dir / "notes" / "total.md"

    if not images_dir.exists():
        print(f"Error: {images_dir} not found")
        sys.exit(1)

    # Collect PNG files sorted by leading number
    png_files = sorted(
        [f for f in images_dir.iterdir() if f.is_file() and f.suffix == ".png"],
        key=lambda p: p.name,
    )

    if not png_files:
        print(f"Error: No PNG files found in {images_dir}")
        sys.exit(1)

    print(f"Found {len(png_files)} slide images")

    # Parse speaker notes
    notes_map = parse_notes(notes_path)
    if notes_map:
        print(f"Loaded speaker notes for {len(notes_map)} slides")
    else:
        print("No speaker notes found (notes/total.md missing or empty)")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)

    notes_embedded = 0
    for png_path in png_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(
            str(png_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Embed speaker notes
        slide_num = png_path.name.split("-")[0]
        if slide_num in notes_map:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_map[slide_num]
            notes_embedded += 1

    # Save
    topic_slug = find_topic_slug(deck_dir)
    pptx_path = deck_dir / f"{topic_slug}.pptx"
    prs.save(str(pptx_path))

    size_mb = pptx_path.stat().st_size / 1024 / 1024
    print(f"\nPPTX saved: {pptx_path} ({size_mb:.1f}MB)")
    print(f"Slides: {len(png_files)}, Speaker notes: {notes_embedded}/{len(png_files)}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 image-export.py <slide-deck-dir>")
        sys.exit(1)

    deck_dir = Path(sys.argv[1])
    if not deck_dir.is_dir():
        print(f"Error: {deck_dir} is not a directory")
        sys.exit(1)

    export_pptx(deck_dir)


if __name__ == "__main__":
    main()
