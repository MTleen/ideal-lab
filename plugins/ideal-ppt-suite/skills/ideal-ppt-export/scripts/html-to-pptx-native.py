#!/usr/bin/env python3
"""
Convert HTML slide files to PPTX using a hybrid renderer:
pixel-perfect Playwright screenshots as slide backgrounds plus editable native
text boxes and table shapes overlaid at the same coordinates.
"""

from __future__ import annotations

import argparse
import math
import re
import sys
import tempfile
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


VIEWPORT_WIDTH = 1280
VIEWPORT_HEIGHT = 720
SLIDE_WIDTH_EMU = 12_192_000
SLIDE_HEIGHT_EMU = 6_858_000
EMU_PER_PX = 9525
SCALE_X = SLIDE_WIDTH_EMU / VIEWPORT_WIDTH
SCALE_Y = SLIDE_HEIGHT_EMU / VIEWPORT_HEIGHT
FONT_FALLBACK = "Microsoft YaHei"
MIN_SCREENSHOT_BYTES = 10 * 1024


TEXT_EXTRACTION_JS = r"""
(() => {
  const skippedTags = new Set([
    'SCRIPT', 'STYLE', 'NOSCRIPT', 'META', 'LINK', 'HEAD', 'TITLE', 'HTML',
    'BR', 'HR', 'SVG', 'PATH', 'IMG', 'TABLE', 'THEAD', 'TBODY', 'TR', 'TH', 'TD'
  ]);

  function normalizeText(text) {
    return text
      .replace(/\u00a0/g, ' ')
      .replace(/[ \t\f\v]+/g, ' ')
      .replace(/[ \t\f\v]*\n[ \t\f\v]*/g, '\n')
      .replace(/\n{3,}/g, '\n\n')
      .trim();
  }

  function ownDirectText(el) {
    const parts = [];
    for (const node of el.childNodes) {
      if (node.nodeType === Node.TEXT_NODE) {
        parts.push(node.textContent || '');
      } else if (node.nodeType === Node.ELEMENT_NODE && node.tagName === 'BR') {
        parts.push('\n');
      }
    }
    return normalizeText(parts.join(''));
  }

  function isVisible(el) {
    for (let current = el; current && current.nodeType === Node.ELEMENT_NODE; current = current.parentElement) {
      const style = window.getComputedStyle(current);
      if (style.display === 'none' || style.visibility === 'hidden' || parseFloat(style.opacity || '1') === 0) {
        return false;
      }
    }
    return true;
  }

  return Array.from(document.querySelectorAll('*')).filter(el => {
    if (skippedTags.has(el.tagName) || el.closest('table')) return false;
    if (!isVisible(el)) return false;
    return ownDirectText(el).length > 0;
  }).map(el => {
  const rect = el.getBoundingClientRect();
  const style = window.getComputedStyle(el);
  return {
    tag: el.tagName,
    text: ownDirectText(el),
    x: rect.x, y: rect.y, width: rect.width, height: rect.height,
    fontSize: parseFloat(style.fontSize),
    fontWeight: style.fontWeight,
    color: style.color,
    fontFamily: style.fontFamily,
    textAlign: style.textAlign,
    lineHeight: parseFloat(style.lineHeight)
  };
  }).filter(item => item.text && item.width > 0 && item.height > 0);
})()
"""


TABLE_EXTRACTION_JS = r"""
Array.from(document.querySelectorAll('table')).map((el, index) => {
  const rect = el.getBoundingClientRect();
  const style = window.getComputedStyle(el);
  return {
    index,
    html: el.outerHTML,
    x: rect.x,
    y: rect.y,
    width: rect.width,
    height: rect.height,
    fontSize: parseFloat(style.fontSize),
    color: style.color
  };
}).filter(item => item.width > 0 && item.height > 0);
"""


BODY_BACKGROUND_JS = r"""
(() => {
  const style = window.getComputedStyle(document.body);
  return {
    background: style.background,
    backgroundColor: style.backgroundColor,
    backgroundImage: style.backgroundImage
  };
})()
"""


HIDE_TEXT_JS = r"""
document.body.insertAdjacentHTML(
  'beforeend',
  '<style id="hide-text">*{color:transparent!important;text-shadow:none!important;}</style>'
);
"""


def emu_x(px: float) -> Emu:
    return Emu(int(round(px * SCALE_X)))


def emu_y(px: float) -> Emu:
    return Emu(int(round(px * SCALE_Y)))


def clamped_rect(item: dict[str, Any]) -> tuple[float, float, float, float] | None:
    left = max(0.0, min(VIEWPORT_WIDTH, float(item.get("x", 0))))
    top = max(0.0, min(VIEWPORT_HEIGHT, float(item.get("y", 0))))
    right = max(left + 1.0, min(VIEWPORT_WIDTH, float(item.get("x", 0)) + max(1.0, float(item.get("width", 1)))))
    bottom = max(top + 1.0, min(VIEWPORT_HEIGHT, float(item.get("y", 0)) + max(1.0, float(item.get("height", 1)))))
    width = right - left
    height = bottom - top
    if width <= 0 or height <= 0:
        return None
    return left, top, width, height


def dedupe_text_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[tuple[str, int, int]] = set()
    occupied_positions: set[tuple[int, int]] = set()
    deduped: list[dict[str, Any]] = []
    for item in items:
        text = re.sub(r"\s+", " ", str(item.get("text", ""))).strip()
        if not text:
            continue
        rect = clamped_rect(item)
        if rect is None:
            continue
        item = dict(item)
        item["text"] = text
        left, top, width, height = rect
        key = (text, round(left), round(top))
        if key in seen:
            continue
        seen.add(key)

        # Fallback layout can assign adjacent labels to the same anchor. Keep
        # all editable text, but move later boxes to the next available row.
        while (round(left), round(top)) in occupied_positions and top + height < VIEWPORT_HEIGHT:
            top = min(VIEWPORT_HEIGHT - height, top + height + 4)
        occupied_positions.add((round(left), round(top)))
        item["x"], item["y"], item["width"], item["height"] = left, top, width, height
        deduped.append(item)
    return deduped


def px_to_pt(px: float | None, default: float = 12.0) -> float:
    if px is None or not math.isfinite(px) or px <= 0:
        return default
    return max(1.0, px * 0.75)


def parse_rgb(css_color: str | None) -> RGBColor | None:
    if not css_color:
        return None

    color = css_color.strip()
    if color.lower() in {"transparent", "currentcolor"}:
        return None

    hex_match = re.search(r"#([0-9a-fA-F]{3,8})", color)
    if hex_match:
        value = hex_match.group(1)
        if len(value) in {3, 4}:
            value = "".join(ch * 2 for ch in value[:3])
        elif len(value) in {6, 8}:
            value = value[:6]
        if len(value) == 6:
            return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    rgb_match = re.search(
        r"rgba?\(\s*([0-9.]+)\s*,\s*([0-9.]+)\s*,\s*([0-9.]+)",
        color,
        flags=re.IGNORECASE,
    )
    if rgb_match:
        channels = [max(0, min(255, int(float(rgb_match.group(i))))) for i in range(1, 4)]
        return RGBColor(channels[0], channels[1], channels[2])

    named = {
        "black": RGBColor(0, 0, 0),
        "white": RGBColor(255, 255, 255),
        "red": RGBColor(255, 0, 0),
        "green": RGBColor(0, 128, 0),
        "blue": RGBColor(0, 0, 255),
        "gray": RGBColor(128, 128, 128),
        "grey": RGBColor(128, 128, 128),
    }
    return named.get(color.lower())


def first_gradient_color(background: dict[str, str]) -> RGBColor | None:
    combined = " ".join(background.get(k, "") for k in ("backgroundImage", "background", "backgroundColor"))
    if "linear-gradient" not in combined:
        return parse_rgb(background.get("backgroundColor") or background.get("background"))
    return parse_rgb(combined)


def is_dark_color(color: RGBColor | None) -> bool:
    if color is None:
        return False
    luminance = (0.2126 * color[0]) + (0.7152 * color[1]) + (0.0722 * color[2])
    return luminance < 80


def text_alignment(value: str | None) -> PP_ALIGN:
    normalized = (value or "").lower()
    if normalized in {"center", "-webkit-center"}:
        return PP_ALIGN.CENTER
    if normalized in {"right", "end", "-webkit-right"}:
        return PP_ALIGN.RIGHT
    if normalized == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def is_bold(weight: str | int | float | None) -> bool:
    if weight is None:
        return False
    try:
        return int(float(str(weight))) >= 600
    except ValueError:
        return str(weight).lower() in {"bold", "bolder"}


def normalized_font_name(_: str | None) -> str:
    return FONT_FALLBACK


def add_background(slide, screenshot_path: Path) -> None:
    slide.shapes.add_picture(
        str(screenshot_path),
        Emu(0),
        Emu(0),
        width=Emu(SLIDE_WIDTH_EMU),
        height=Emu(SLIDE_HEIGHT_EMU),
    )


def add_text_overlay(slide, item: dict[str, Any], dark_background: bool) -> None:
    rect = clamped_rect(item)
    if rect is None:
        return
    left, top, width, height = rect

    textbox = slide.shapes.add_textbox(emu_x(left), emu_y(top), emu_x(width), emu_y(height))
    textbox.fill.background()
    textbox.line.fill.background()

    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = Emu(0)
    frame.margin_right = Emu(0)
    frame.margin_top = Emu(0)
    frame.margin_bottom = Emu(0)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    paragraph = frame.paragraphs[0]
    paragraph.alignment = text_alignment(item.get("textAlign"))
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.0

    run = paragraph.add_run()
    run.text = item["text"]
    font = run.font
    font.name = normalized_font_name(item.get("fontFamily"))
    font.size = Pt(px_to_pt(item.get("fontSize")))
    font.bold = is_bold(item.get("fontWeight"))

    color = parse_rgb(item.get("color"))
    if color is None and dark_background:
        color = RGBColor(255, 255, 255)
    if color is not None:
        font.color.rgb = color


def table_rows_from_html(html: str) -> list[list[str]]:
    soup = BeautifulSoup(html, "html.parser")
    rows: list[list[str]] = []
    for tr in soup.find_all("tr"):
        cells = [cell.get_text(" ", strip=True) for cell in tr.find_all(["th", "td"])]
        if cells:
            rows.append(cells)
    if not rows:
        return []

    max_cols = max(len(row) for row in rows)
    return [row + [""] * (max_cols - len(row)) for row in rows]


def set_cell_text(cell, text: str, row_index: int) -> None:
    cell.text = text
    cell.margin_left = Emu(4 * EMU_PER_PX)
    cell.margin_right = Emu(4 * EMU_PER_PX)
    cell.margin_top = Emu(2 * EMU_PER_PX)
    cell.margin_bottom = Emu(2 * EMU_PER_PX)

    paragraph = cell.text_frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.name = FONT_FALLBACK
        run.font.size = Pt(9)
        run.font.bold = row_index == 0
        run.font.color.rgb = RGBColor(17, 24, 39) if row_index == 0 else RGBColor(75, 85, 99)


def add_table_overlay(slide, table_info: dict[str, Any]) -> None:
    rows = table_rows_from_html(table_info["html"])
    if not rows:
        return
    rect = clamped_rect(table_info)
    if rect is None:
        return
    left, top, width, height = rect

    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    shape = slide.shapes.add_table(
        row_count,
        col_count,
        emu_x(left),
        emu_y(top),
        emu_x(width),
        emu_y(height),
    )
    table = shape.table

    for i in range(row_count):
        table.rows[i].height = Emu(int(round(height * SCALE_Y / row_count)))
    for j in range(col_count):
        table.columns[j].width = Emu(int(round(width * SCALE_X / col_count)))

    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            set_cell_text(cell, value, i)
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = RGBColor(243, 244, 246)
            elif i % 2:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(249, 250, 251)


def slide_files(html_dir: Path) -> list[Path]:
    files = list(html_dir.glob("*-slide-*.html"))
    if not files:
        files = list(html_dir.glob("*.html"))

    def sort_key(path: Path) -> tuple[int, str]:
        match = re.match(r"(\d+)", path.name)
        return (int(match.group(1)) if match else 10**9, path.name)

    return sorted(files, key=sort_key)


def render_html_to_slide(page, prs: Presentation, html_path: Path, screenshot_dir: Path) -> None:
    page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
    page.emulate_media(media="screen")
    page.wait_for_load_state("networkidle")
    page.wait_for_timeout(2000)

    text_items = page.evaluate(TEXT_EXTRACTION_JS)
    table_items = page.evaluate(TABLE_EXTRACTION_JS)
    body_background = page.evaluate(BODY_BACKGROUND_JS)
    dark_background = is_dark_color(first_gradient_color(body_background))

    page.evaluate(HIDE_TEXT_JS)
    page.wait_for_timeout(500)

    screenshot_path = screenshot_dir / f"{html_path.stem}.png"
    page.screenshot(path=str(screenshot_path), full_page=False)
    screenshot_size = screenshot_path.stat().st_size
    if screenshot_size < MIN_SCREENSHOT_BYTES:
        raise RuntimeError(
            f"Playwright rendering failed for {html_path.name}: "
            f"screenshot is only {screenshot_size} bytes"
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, screenshot_path)
    for item in dedupe_text_items(text_items):
        add_text_overlay(slide, item, dark_background)
    for table_info in table_items:
        add_table_overlay(slide, table_info)


def convert(html_dir: Path, output_path: Path, device_scale: float) -> None:
    files = slide_files(html_dir)
    if not files:
        raise FileNotFoundError(f"No HTML files found in {html_dir}")

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    with tempfile.TemporaryDirectory(prefix="html-to-pptx-native-") as temp_dir:
        screenshot_dir = Path(temp_dir)
        try:
            with sync_playwright() as playwright:
                browser = launch_browser(playwright)
                context = browser.new_context(
                    viewport={"width": VIEWPORT_WIDTH, "height": VIEWPORT_HEIGHT},
                    device_scale_factor=device_scale,
                )
                page = context.new_page()
                for html_path in files:
                    print(f"Rendering {html_path.name}", file=sys.stderr)
                    render_html_to_slide(page, prs, html_path, screenshot_dir)
                context.close()
                browser.close()
        except Exception as exc:
            print(f"WARNING: Playwright rendering failed ({exc}). Using static fallback renderer.", file=sys.stderr)
            prs = Presentation()
            prs.slide_width = Emu(SLIDE_WIDTH_EMU)
            prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
            for html_path in files:
                print(f"Fallback rendering {html_path.name}", file=sys.stderr)
                render_html_fallback_to_slide(prs, html_path, screenshot_dir)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)


def launch_browser(playwright):
    launch_options: list[dict[str, Any]] = [{}]
    for executable in (
        Path("/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"),
        Path("/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge"),
    ):
        if executable.exists():
            launch_options.append({"executable_path": str(executable)})

    last_error: Exception | None = None
    for options in launch_options:
        try:
            return playwright.chromium.launch(
                headless=True,
                args=["--no-sandbox", "--disable-setuid-sandbox", "--single-process"],
                **options,
            )
        except Exception as exc:
            last_error = exc
    assert last_error is not None
    raise last_error


def extract_css_vars(html: str) -> dict[str, str]:
    vars_map: dict[str, str] = {}
    for match in re.finditer(r"--([\w-]+)\s*:\s*([^;]+);", html):
        vars_map[match.group(1)] = match.group(2).strip()
    return vars_map


def resolve_css_value(value: str | None, css_vars: dict[str, str]) -> str | None:
    if not value:
        return value
    match = re.search(r"var\(--([\w-]+)\)", value)
    if match:
        return css_vars.get(match.group(1), value)
    return value


def body_background_from_html(html: str, css_vars: dict[str, str]) -> RGBColor:
    body_rule = re.search(r"body\s*\{([^}]+)\}", html, flags=re.DOTALL)
    if not body_rule:
        return RGBColor(255, 255, 255)
    declarations = body_rule.group(1)
    bg_match = re.search(r"background\s*:\s*([^;]+)", declarations)
    bg_value = resolve_css_value(bg_match.group(1).strip(), css_vars) if bg_match else "#FFFFFF"
    return parse_rgb(bg_value) or RGBColor(255, 255, 255)


def save_fallback_background(path: Path, background: RGBColor, dark: bool) -> None:
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (VIEWPORT_WIDTH, VIEWPORT_HEIGHT), tuple(background))
    draw = ImageDraw.Draw(image)
    if dark:
        for y in range(VIEWPORT_HEIGHT):
            ratio = y / max(1, VIEWPORT_HEIGHT - 1)
            r = int(background[0] * (1 - ratio) + 15 * ratio)
            g = int(background[1] * (1 - ratio) + 23 * ratio)
            b = int(background[2] * (1 - ratio) + 42 * ratio)
            draw.line([(0, y), (VIEWPORT_WIDTH, y)], fill=(r, g, b))
    image.save(path)


def fallback_color_for_element(element, css_vars: dict[str, str], dark: bool) -> str:
    style = element.get("style", "")
    inline = re.search(r"color\s*:\s*([^;]+)", style)
    if inline:
        return resolve_css_value(inline.group(1).strip(), css_vars) or "rgb(255,255,255)"

    classes = set(element.get("class", []))
    class_color = {
        "accent-blue": "--mod-blue",
        "accent-green": "--mod-green",
        "accent-red": "--mod-red",
        "accent-purple": "--mod-purple",
        "accent-amber": "--mod-amber",
        "n1": "--mod-blue",
        "n2": "--mod-green",
        "n3": "--mod-amber",
        "n4": "--mod-red",
        "n5": "--mod-purple",
    }
    for cls, var_name in class_color.items():
        if cls in classes:
            return css_vars.get(var_name[2:], "#111827")
    if any(cls in classes for cls in ("kpi-value", "stat-value")):
        return css_vars.get("mod-amber", "#D97706")
    if dark:
        return "rgb(255,255,255)"
    if any(cls in classes for cls in ("card-body", "kpi-label", "stat-label")):
        return css_vars.get("text-secondary", "#4B5563")
    return css_vars.get("text", "#111827")


def fallback_font_size(element) -> float:
    classes = set(element.get("class", []))
    if "title" in classes:
        return 44
    if "page-title" in classes:
        return 24
    if "subtitle" in classes:
        return 18
    if "kpi-value" in classes or "stat-value" in classes:
        return 28
    if "card-title" in classes or "section-title" in classes or "side-title" in classes:
        return 15
    if element.name in {"h1"}:
        return 36
    if element.name in {"h2", "h3"}:
        return 22
    if element.name == "li":
        return 12
    return 13


def fallback_text_elements(soup: BeautifulSoup) -> list[Any]:
    blocked = {"head", "title", "meta", "link", "script", "style", "noscript", "table", "thead", "tbody", "tr", "th", "td"}
    elements = []
    for element in soup.find_all(True):
        if element.name in blocked or element.find_parent("table"):
            continue
        has_direct_text = any(
            getattr(child, "name", None) is None and str(child).strip()
            for child in element.children
        )
        if has_direct_text and element.get_text(" ", strip=True):
            elements.append(element)
    return elements


def fallback_layout_item(index: int, element, slide_number: int, css_vars: dict[str, str], dark: bool) -> dict[str, Any]:
    classes = set(element.get("class", []))
    text = element.get_text(" ", strip=True)
    font_size = fallback_font_size(element)
    width = 520
    height = max(20, font_size * 1.35)
    x = 50
    y = 40 + index * 28

    if slide_number in {1, 5}:
        if "badge" in classes:
            x, y, width = 60, 180, 520
        elif "title" in classes or "page-title" in classes:
            x, y, width, height = 60, 220 if slide_number == 1 else 45, 720, 80
        elif "subtitle" in classes:
            x, y, width, height = 60, 335, 720, 60
        elif "meta" in classes:
            x, y, width = 60, 430, 900
        elif "takeaway-num" in classes:
            n = max(0, len([e for e in element.find_all_previous(class_="takeaway-num")]))
            x, y, width = 60, 120 + n * 82, 45
        elif "takeaway-text" in classes:
            n = max(0, len([e for e in element.find_all_previous(class_="takeaway-text")]))
            x, y, width, height = 115, 115 + n * 82, 570, 58
        elif "kpi-value" in classes or "kpi-label" in classes:
            n = len(element.find_all_previous(class_=element.get("class", [""])[0]))
            x = 760 + (n % 2) * 210
            y = 150 + (n // 2) * 110 + (42 if "kpi-label" in classes else 0)
            width = 180
        elif "stat-value" in classes or "stat-label" in classes:
            n = len(element.find_all_previous(class_=element.get("class", [""])[0]))
            x, y, width = 860, 190 + n * 76 + (36 if "stat-label" in classes else 0), 260
        elif "highlight-item" in classes:
            n = len(element.find_all_previous(class_="highlight-item"))
            x, y, width = 860, 500 + n * 28, 360
        elif "cta-title" in classes or "cta-items" in classes or element.name == "li":
            n = len(element.find_all_previous(["li"]))
            x, y, width = 755, 410 + n * 30, 410
    else:
        if "page-title" in classes:
            x, y, width, height = 40, 28, 900, 34
        elif "card-title" in classes:
            n = len(element.find_all_previous(class_="card-title"))
            if slide_number == 3:
                x, y, width = 760, 75 + n * 128, 390
            elif slide_number == 4 and n >= 2:
                x, y, width = 55 + (n - 2) * 405, 555, 360
            else:
                x, y, width = 60 + (n % 2) * 600, 90 + (n // 2) * 145, 520
        elif "card-body" in classes:
            n = len(element.find_all_previous(class_="card-body"))
            if slide_number == 3:
                x, y, width, height = 760, 102 + n * 128, 400, 74
            elif slide_number == 4 and n >= 0:
                x, y, width, height = 55 + n * 405, 582, 360, 55
            else:
                x, y, width, height = 60 + (n % 2) * 600, 118 + (n // 2) * 145, 520, 72
        elif element.name == "li":
            n = len(element.find_all_previous("li"))
            x, y, width = 80 + ((n // 4) % 2) * 600, 118 + (n % 4) * 22 + (n // 8) * 145, 500
        elif "layer-label" in classes:
            n = len(element.find_all_previous(class_="layer-label"))
            x, y, width = 60, 82 + n * 92, 190
        elif "arch-tag" in classes:
            n = len(element.find_all_previous(class_="arch-tag"))
            x, y, width = 250 + (n % 3) * 145, 82 + (n // 5) * 92 + ((n % 5) // 3) * 26, 135
        elif "kpi-value" in classes or "kpi-label" in classes:
            n = len(element.find_all_previous(class_=element.get("class", [""])[0]))
            x = 60 + (n % 4) * 290 if slide_number == 2 else 755 + (n % 2) * 210
            y = 590 if slide_number == 2 else 500 + (n // 2) * 80
            y += 34 if "kpi-label" in classes else 0
            width = 180
        elif "section-title" in classes:
            x, y, width = 40, 245, 400
        elif "side-title" in classes or "side-list" in classes:
            n = len(element.find_all_previous(class_=element.get("class", [""])[0]))
            x, y, width = 60 + (n % 2) * 600, 90 + ("side-list" in classes) * 32, 520

    return {
        "tag": element.name.upper(),
        "text": text,
        "x": x,
        "y": y,
        "width": width,
        "height": height,
        "fontSize": font_size,
        "fontWeight": "700" if font_size >= 15 or element.name in {"strong", "b"} else "400",
        "color": fallback_color_for_element(element, css_vars, dark),
        "fontFamily": FONT_FALLBACK,
        "textAlign": "left",
        "lineHeight": font_size * 1.3,
    }


def fallback_table_infos(soup: BeautifulSoup, slide_number: int) -> list[dict[str, Any]]:
    infos = []
    for index, table in enumerate(soup.find_all("table")):
        infos.append(
            {
                "index": index,
                "html": str(table),
                "x": 40,
                "y": 285 if slide_number == 4 else 180,
                "width": 1200,
                "height": 170 if slide_number == 4 else 260,
                "fontSize": 12,
                "color": "#111827",
            }
        )
    return infos


def render_html_fallback_to_slide(prs: Presentation, html_path: Path, screenshot_dir: Path) -> None:
    html = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")
    css_vars = extract_css_vars(html)
    background = body_background_from_html(html, css_vars)
    dark = is_dark_color(background)

    screenshot_path = screenshot_dir / f"{html_path.stem}-fallback.png"
    save_fallback_background(screenshot_path, background, dark)

    match = re.match(r"(\d+)", html_path.name)
    slide_number = int(match.group(1)) if match else 0
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, screenshot_path)

    fallback_items = [
        fallback_layout_item(index, element, slide_number, css_vars, dark)
        for index, element in enumerate(fallback_text_elements(soup))
    ]
    for item in dedupe_text_items(fallback_items):
        add_text_overlay(slide, item, dark)
    for table_info in fallback_table_infos(soup, slide_number):
        add_table_overlay(slide, table_info)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert HTML slides to hybrid native PPTX.")
    parser.add_argument("html_dir", type=Path, help="Directory containing NN-slide-*.html files.")
    parser.add_argument("-o", "--output", required=True, type=Path, help="Output PPTX path.")
    parser.add_argument("--device-scale", default=2, type=float, help="Playwright device scale factor.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        convert(args.html_dir, args.output, args.device_scale)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    print(f"Wrote {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
