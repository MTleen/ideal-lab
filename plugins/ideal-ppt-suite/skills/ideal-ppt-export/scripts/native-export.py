#!/usr/bin/env python3
"""
Native Shapes PPTX Export - using ppt-master's svg_to_shapes converter
Converts SVG elements to DrawingML shapes for fully editable PPTX.

Usage: python3 native-export.py <slide-deck-dir>
"""

import sys
import os
import re
import shutil
import zipfile
import tempfile
from pathlib import Path
from xml.etree import ElementTree as ET

try:
    from pptx import Presentation
except ImportError:
    print("Error: pip install python-pptx"); sys.exit(1)

# Add this script's directory to path (svg_to_shapes.py is in the same dir)
script_dir = Path(__file__).parent
sys.path.insert(0, str(script_dir))

try:
    from svg_to_shapes import convert_svg_to_slide_shapes
except ImportError:
    print("Error: svg_to_shapes.py not found in pptmaster_scripts/"); sys.exit(1)

EMU_PER_INCH = 914400
EMU_PER_PX = EMU_PER_INCH / 96

NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}
for prefix, uri in NAMESPACES.items():
    ET.register_namespace(prefix, uri)


def create_notes_slide_xml(notes_text):
    notes_escaped = notes_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    paragraphs = []
    for para in notes_escaped.split('\n'):
        if para.strip():
            paragraphs.append(f'<a:p><a:r><a:rPr lang="zh-CN" dirty="0"/><a:t>{para}</a:t></a:r></a:p>')
        else:
            paragraphs.append('<a:p><a:endParaRPr lang="zh-CN" dirty="0"/></a:p>')
    paras_xml = '\n'.join(paragraphs)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
         xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>{paras_xml}</p:txBody></p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 native-export.py <slide-deck-dir>"); sys.exit(1)

    project_dir = Path(sys.argv[1]).resolve()
    slug = project_dir.name

    svg_dir = project_dir / 'svg_final' if (project_dir / 'svg_final').exists() else project_dir / 'svg_output'
    if not svg_dir.exists():
        print(f"Error: No SVG directory"); sys.exit(1)

    svg_files = sorted(svg_dir.glob('*.svg'), key=lambda f: int(re.match(r'(\d+)', f.name).group(1)))
    if not svg_files:
        print(f"Error: No SVG files"); sys.exit(1)

    print(f"Found {len(svg_files)} slides")

    svg_w, svg_h = 1280, 720
    width_emu = int(svg_w * EMU_PER_PX)
    height_emu = int(svg_h * EMU_PER_PX)

    temp_dir = Path(tempfile.mkdtemp())
    try:
        # 1. Create base PPTX
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)
        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # 2. Extract
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)
        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        # 3. Load notes
        notes_dir = project_dir / 'notes'
        notes_map = {}
        if notes_dir.exists():
            for nf in notes_dir.glob('*.md'):
                num_match = re.match(r'(\d+)', nf.name)
                if num_match:
                    notes_map[int(num_match.group(1))] = nf.read_text('utf-8').strip()

        # 4. Process each SVG with native shapes conversion
        for i, svg_path in enumerate(svg_files, 1):
            try:
                slide_xml, media_files, rel_entries = convert_svg_to_slide_shapes(
                    svg_path, slide_num=i, verbose=False
                )

                # Write slide XML
                slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{i}.xml'
                slide_xml_path.write_text(slide_xml, encoding='utf-8')

                # Write media files
                for media_name, media_data in media_files.items():
                    with open(media_dir / media_name, 'wb') as f:
                        f.write(media_data)

                # Build relationships XML
                rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                rels_dir.mkdir(exist_ok=True)
                rels_path = rels_dir / f'slide{i}.xml.rels'

                extra_rels = ''
                for rel in rel_entries:
                    extra_rels += f'\n  <Relationship Id="{rel["id"]}" Type="{rel["type"]}" Target="{rel["target"]}"/>'

                rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                rels_path.write_text(rels_xml, encoding='utf-8')

                # Notes
                if i in notes_map and notes_map[i]:
                    notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                    notes_slides_dir.mkdir(exist_ok=True)
                    (notes_slides_dir / f'notesSlide{i}.xml').write_text(
                        create_notes_slide_xml(notes_map[i]), encoding='utf-8'
                    )
                    notes_rels_dir = notes_slides_dir / '_rels'
                    notes_rels_dir.mkdir(exist_ok=True)
                    (notes_rels_dir / f'notesSlide{i}.xml.rels').write_text(
                        f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{i}.xml"/>
</Relationships>''', encoding='utf-8'
                    )
                    # Add notes rel to slide rels
                    rels_content = rels_path.read_text('utf-8')
                    notes_rel = f'  <Relationship Id="rId10" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{i}.xml"/>'
                    rels_content = rels_content.replace('</Relationships>', notes_rel + '\n</Relationships>')
                    rels_path.write_text(rels_content, encoding='utf-8')

                print(f"  [{i}/{len(svg_files)}] {svg_path.name}" + (" +notes" if i in notes_map else ""))

            except Exception as e:
                print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")

        # 5. Update [Content_Types].xml
        ct_path = extract_dir / '[Content_Types].xml'
        ct = ct_path.read_text('utf-8')
        additions = []
        if 'Extension="png"' not in ct:
            additions.append('<Default Extension="png" ContentType="image/png"/>')
        if 'Extension="jpg"' not in ct:
            additions.append('<Default Extension="jpg" ContentType="image/jpeg"/>')
        for i in range(1, len(svg_files) + 1):
            if i in notes_map:
                override = f'<Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                if override not in ct:
                    additions.append(override)
        if additions:
            ct = ct.replace('</Types>', '\n'.join(additions) + '\n</Types>')
            ct_path.write_text(ct, encoding='utf-8')

        # 6. Repackage
        output_path = project_dir / f'{slug}.pptx'
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for fp in extract_dir.rglob('*'):
                if fp.is_file():
                    zf.write(fp, fp.relative_to(extract_dir))

        print(f"\nExported: {output_path} ({len(svg_files)} slides, native shapes)")

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


if __name__ == '__main__':
    main()
