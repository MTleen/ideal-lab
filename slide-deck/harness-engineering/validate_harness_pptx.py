#!/usr/bin/env python3
from __future__ import annotations

import argparse
import io
import math
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from xml.etree import ElementTree as ET

from PIL import Image, ImageChops, ImageStat
from pptx import Presentation

try:
    from skimage.metrics import structural_similarity as ssim
except Exception:  # pragma: no cover - optional dependency
    ssim = None


BASE_DIR = Path("/Users/mathrippermacmini/Documents/Sync/Work/电信/理想/产品/0-最佳实践/slide-deck/harness-engineering")
DEFAULT_PPTX = BASE_DIR / "harness-engineering.pptx"
DEFAULT_SCREENSHOTS = BASE_DIR / "screenshots"
DEFAULT_DEBUG = BASE_DIR / "pptx_debug"

EXPECTED_SLIDES = 5
EXPECTED_BG_SIZE = (2560, 1440)
MIN_BG_BYTES = 10 * 1024
SLIDE_WIDTH_EMU = 12_192_000
SLIDE_HEIGHT_EMU = 6_858_000
SIMILARITY_THRESHOLD = 0.94

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


@dataclass
class BackgroundResult:
    slide: int
    status: str
    dimensions: tuple[int, int] | None
    size_kb: float
    path: Path | None
    reason: str


@dataclass
class SimilarityResult:
    slide: int
    status: str
    score: float | None
    reason: str


@dataclass
class TextResult:
    slide: int
    textbox_count: int
    coords_status: str
    title_leaks: int
    failures: list[str]


def parse_xml(data: bytes) -> ET.Element:
    return ET.fromstring(data)


def rel_targets(zf: zipfile.ZipFile, slide_index: int) -> dict[str, str]:
    rels_name = f"ppt/slides/_rels/slide{slide_index}.xml.rels"
    if rels_name not in zf.namelist():
        return {}
    root = parse_xml(zf.read(rels_name))
    result = {}
    for rel in root.findall("rel:Relationship", NS):
        rid = rel.attrib.get("Id")
        target = rel.attrib.get("Target", "")
        if rid:
            result[rid] = target
    return result


def resolve_slide_target(slide_index: int, target: str) -> str:
    slide_dir = PurePosixPath(f"ppt/slides/slide{slide_index}.xml").parent
    path = PurePosixPath(target)
    if not path.is_absolute():
        path = slide_dir / path
    parts: list[str] = []
    for part in path.parts:
        if part in ("", "."):
            continue
        if part == "..":
            if parts:
                parts.pop()
        else:
            parts.append(part)
    return "/".join(parts)


def image_from_rid(zf: zipfile.ZipFile, rels: dict[str, str], slide_index: int, rid: str) -> Image.Image | None:
    target = rels.get(rid)
    if not target:
        return None
    image_name = resolve_slide_target(slide_index, target)
    if image_name not in zf.namelist():
        return None
    return Image.open(io.BytesIO(zf.read(image_name))).convert("RGB")


def shape_name(node: ET.Element) -> str:
    for tag in ("p:nvPicPr/p:cNvPr", "p:nvSpPr/p:cNvPr"):
        c_nv_pr = node.find(tag, NS)
        if c_nv_pr is not None:
            return c_nv_pr.attrib.get("name", "")
    return ""


def xfrm_of(node: ET.Element) -> tuple[int, int, int, int] | None:
    xfrm = node.find(".//a:xfrm", NS)
    if xfrm is None:
        return None
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return None
    return (
        int(off.attrib.get("x", "0")),
        int(off.attrib.get("y", "0")),
        int(ext.attrib.get("cx", "0")),
        int(ext.attrib.get("cy", "0")),
    )


def is_full_slide(box: tuple[int, int, int, int] | None) -> bool:
    if box is None:
        return False
    x, y, cx, cy = box
    return abs(x) <= 2 and abs(y) <= 2 and abs(cx - SLIDE_WIDTH_EMU) <= 2 and abs(cy - SLIDE_HEIGHT_EMU) <= 2


def extract_background_candidates(zf: zipfile.ZipFile, slide_index: int, root: ET.Element) -> list[Image.Image]:
    rels = rel_targets(zf, slide_index)
    candidates: list[tuple[int, Image.Image]] = []

    bg = root.find("p:cSld/p:bg", NS)
    if bg is not None:
        blip = bg.find(".//a:blip", NS)
        rid = blip.attrib.get(f"{{{NS['r']}}}embed") if blip is not None else None
        if rid:
            image = image_from_rid(zf, rels, slide_index, rid)
            if image is not None:
                candidates.append((0, image))

    for pic in root.findall(".//p:pic", NS):
        blip = pic.find(".//a:blip", NS)
        rid = blip.attrib.get(f"{{{NS['r']}}}embed") if blip is not None else None
        if not rid:
            continue
        name = shape_name(pic).lower()
        box = xfrm_of(pic)
        priority = 10
        if "background" in name:
            priority = 1
        elif is_full_slide(box):
            priority = 2
        else:
            continue
        image = image_from_rid(zf, rels, slide_index, rid)
        if image is not None:
            candidates.append((priority, image))

    candidates.sort(key=lambda item: item[0])
    return [image for _, image in candidates]


def visible_content(image: Image.Image) -> bool:
    gray = image.convert("L")
    stat = ImageStat.Stat(gray)
    if not stat.var or stat.var[0] < 2.0:
        return False
    colors = image.resize((64, 36), Image.Resampling.BILINEAR).getcolors(maxcolors=64 * 36)
    return colors is None or len(colors) > 8


def save_png(image: Image.Image, path: Path) -> int:
    image.save(path, format="PNG")
    return path.stat().st_size


def validate_backgrounds(pptx_path: Path, debug_dir: Path) -> list[BackgroundResult]:
    debug_dir.mkdir(parents=True, exist_ok=True)
    results: list[BackgroundResult] = []
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_index in range(1, EXPECTED_SLIDES + 1):
            slide_name = f"ppt/slides/slide{slide_index}.xml"
            out_path = debug_dir / f"slide_{slide_index}_bg.png"
            if slide_name not in zf.namelist():
                results.append(BackgroundResult(slide_index, "FAIL", None, 0, None, "slide XML missing"))
                continue
            root = parse_xml(zf.read(slide_name))
            candidates = extract_background_candidates(zf, slide_index, root)
            if not candidates:
                results.append(BackgroundResult(slide_index, "FAIL", None, 0, None, "background image not found"))
                continue

            image = candidates[0]
            file_size = save_png(image, out_path)
            dimensions_ok = image.size == EXPECTED_BG_SIZE
            size_ok = file_size > MIN_BG_BYTES
            visible_ok = visible_content(image)
            failures = []
            if not dimensions_ok:
                failures.append(f"expected {EXPECTED_BG_SIZE[0]}x{EXPECTED_BG_SIZE[1]}")
            if not size_ok:
                failures.append("file size <= 10KB")
            if not visible_ok:
                failures.append("blank or solid color")
            status = "PASS" if not failures else "FAIL"
            results.append(
                BackgroundResult(
                    slide_index,
                    status,
                    image.size,
                    file_size / 1024,
                    out_path,
                    "; ".join(failures),
                )
            )
    return results


def screenshot_paths(screenshots_dir: Path) -> list[Path]:
    return sorted(screenshots_dir.glob("*.png"))


def similarity_score(bg_path: Path, screenshot_path: Path) -> float:
    bg = Image.open(bg_path).convert("RGB").resize((1280, 720), Image.Resampling.LANCZOS)
    shot = Image.open(screenshot_path).convert("RGB").resize((1280, 720), Image.Resampling.LANCZOS)
    if ssim is not None:
        import numpy as np

        return float(ssim(np.array(bg), np.array(shot), channel_axis=2, data_range=255))
    diff = ImageChops.difference(bg, shot)
    stat = ImageStat.Stat(diff)
    mae = sum(stat.mean) / (3 * 255)
    return max(0.0, 1.0 - mae)


def validate_similarity(backgrounds: list[BackgroundResult], screenshots_dir: Path) -> list[SimilarityResult]:
    shots = screenshot_paths(screenshots_dir)
    results: list[SimilarityResult] = []
    for bg in backgrounds:
        if bg.path is None or bg.status == "FAIL":
            results.append(SimilarityResult(bg.slide, "FAIL", None, "background unavailable"))
            continue
        if bg.slide > len(shots):
            results.append(SimilarityResult(bg.slide, "FAIL", None, "screenshot missing"))
            continue
        score = similarity_score(bg.path, shots[bg.slide - 1])
        status = "PASS" if score >= SIMILARITY_THRESHOLD else "FAIL"
        reason = "" if status == "PASS" else f"below threshold {SIMILARITY_THRESHOLD:.2f}"
        results.append(SimilarityResult(bg.slide, status, score, reason))
    return results


def text_of(node: ET.Element) -> str:
    return "".join(t.text or "" for t in node.findall(".//a:t", NS)).strip()


def textboxes_in_range(root: ET.Element) -> tuple[int, str, int, list[str]]:
    failures: list[str] = []
    title_leaks = 0
    textbox_count = 0
    for shape in root.findall(".//p:sp", NS):
        tx_body = shape.find("p:txBody", NS)
        if tx_body is None:
            continue
        textbox_count += 1
        text = re.sub(r"\s+", " ", text_of(tx_body))
        if re.search(r"\bSlide\s+\d+\s+-", text, re.IGNORECASE):
            title_leaks += 1
        box = xfrm_of(shape)
        if box is None:
            failures.append("TextBox missing coordinates")
            continue
        x, y, cx, cy = box
        if x < 0 or y < 0 or cx < 0 or cy < 0 or x + cx > SLIDE_WIDTH_EMU or y + cy > SLIDE_HEIGHT_EMU:
            failures.append(f"TextBox out of range: x={x}, y={y}, cx={cx}, cy={cy}")
    coords_status = "PASS" if not failures else "FAIL"
    return textbox_count, coords_status, title_leaks, failures


def validate_textboxes_and_tables(pptx_path: Path) -> tuple[list[TextResult], list[str]]:
    # python-pptx is intentionally loaded here to satisfy the parser requirement and verify the package can read the deck.
    prs = Presentation(str(pptx_path))
    failures: list[str] = []
    if len(prs.slides) != EXPECTED_SLIDES:
        failures.append(f"expected {EXPECTED_SLIDES} slides, found {len(prs.slides)}")
    if prs.slide_width != SLIDE_WIDTH_EMU or prs.slide_height != SLIDE_HEIGHT_EMU:
        failures.append(f"slide size is {prs.slide_width}x{prs.slide_height} EMU")

    results: list[TextResult] = []
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_index in range(1, EXPECTED_SLIDES + 1):
            slide_name = f"ppt/slides/slide{slide_index}.xml"
            if slide_name not in zf.namelist():
                results.append(TextResult(slide_index, 0, "FAIL", 0, ["slide XML missing"]))
                continue
            root = parse_xml(zf.read(slide_name))
            textbox_count, coords_status, title_leaks, slide_failures = textboxes_in_range(root)

            native_tables = len(root.findall(".//a:tbl", NS))
            raster_table_pics = []
            for pic in root.findall(".//p:pic", NS):
                name = shape_name(pic).lower()
                if "table" in name or "表" in name:
                    raster_table_pics.append(shape_name(pic))
            if raster_table_pics and native_tables == 0:
                slide_failures.append(f"possible rasterized table image(s): {', '.join(raster_table_pics)}")

            if title_leaks:
                slide_failures.append(f"{title_leaks} leaked title label(s)")
            results.append(TextResult(slide_index, textbox_count, coords_status, title_leaks, slide_failures))
    return results, failures


def format_dimensions(dimensions: tuple[int, int] | None) -> str:
    if dimensions is None:
        return "missing"
    return f"{dimensions[0]}x{dimensions[1]}"


def format_score(score: float | None) -> str:
    if score is None or math.isnan(score):
        return "N/A"
    return f"{score:.4f}"


def build_report(
    backgrounds: list[BackgroundResult],
    similarities: list[SimilarityResult],
    text_results: list[TextResult],
    global_text_failures: list[str],
) -> tuple[str, bool]:
    failures: list[str] = []
    for bg in backgrounds:
        if bg.status == "FAIL":
            failures.append(f"C1 Slide {bg.slide}: {bg.reason}")
    for sim in similarities:
        if sim.status == "FAIL":
            failures.append(f"C2 Slide {sim.slide}: {sim.reason}")
    for failure in global_text_failures:
        failures.append(f"C3: {failure}")
    for tr in text_results:
        if tr.coords_status == "FAIL" or tr.title_leaks:
            failures.extend(f"C3 Slide {tr.slide}: {failure}" for failure in tr.failures)

    lines = [
        "## 验收报告 (ppt-v4-render-fix)",
        "",
        "### C1: Background image quality",
    ]
    for bg in backgrounds:
        detail = bg.reason if bg.reason else f"{format_dimensions(bg.dimensions)}, {bg.size_kb:.1f} KB"
        lines.append(f"- Slide {bg.slide}: {bg.status} - {detail}")

    lines.extend(["", "### C2: PPTX background vs HTML screenshot similarity"])
    for sim in similarities:
        lines.append(f"- Slide {sim.slide}: {format_score(sim.score)} {sim.status}")

    lines.extend(["", "### C3: TextBox editability and alignment"])
    for tr in text_results:
        lines.append(
            f"- Slide {tr.slide}: {tr.textbox_count} TextBoxes, "
            f"coords in range: {tr.coords_status}, title leaks: {tr.title_leaks}"
        )

    lines.extend(["", "### 总结"])
    if failures:
        lines.append(f"- 需要修复: {'; '.join(failures)}")
    else:
        lines.append("- ALL PASS")
    return "\n".join(lines), not failures


def validate(pptx_path: Path, screenshots_dir: Path, debug_dir: Path) -> tuple[str, bool]:
    backgrounds = validate_backgrounds(pptx_path, debug_dir)
    similarities = validate_similarity(backgrounds, screenshots_dir)
    text_results, global_text_failures = validate_textboxes_and_tables(pptx_path)
    return build_report(backgrounds, similarities, text_results, global_text_failures)


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate harness-engineering PPTX render quality.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX)
    parser.add_argument("--screenshots-dir", type=Path, default=DEFAULT_SCREENSHOTS)
    parser.add_argument("--debug-dir", type=Path, default=DEFAULT_DEBUG)
    args = parser.parse_args()

    report, passed = validate(args.pptx, args.screenshots_dir, args.debug_dir)
    print(report)
    return 0 if passed else 1


if __name__ == "__main__":
    raise SystemExit(main())
